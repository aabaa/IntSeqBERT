"""
Fig.1 IntSeqBERT アーキテクチャ図 — PowerPoint 版
CICM 2026 paper — Figure 1 (editable PPTX)

python-pptx を使って編集可能な PPTX を生成。
座標系: matplotlib の (0–16, 0–7) データ座標を
        スライド幅 33.87cm, 高さ 14.82cm に変換。

出力: paper/cicm2026/figures/fig1_architecture.pptx
"""

from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

OUT_DIR = Path(__file__).resolve().parent.parent / "figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── スライドサイズ ────────────────────────────────────────────────────────
SLIDE_W_CM = 33.87
SLIDE_H_CM = 14.82
DATA_W = 16.0
DATA_H = 7.0


def cm(v):
    return Emu(int(v * 360000))


def to_emu_x(data_x):
    return cm(data_x / DATA_W * SLIDE_W_CM)


def to_emu_y(data_y):
    """data_y → EMU (y反転: matplotlib 下=0 → pptx 上=0)"""
    return cm((DATA_H - data_y) / DATA_H * SLIDE_H_CM)


def to_emu_w(data_w):
    return cm(data_w / DATA_W * SLIDE_W_CM)


def to_emu_h(data_h):
    return cm(data_h / DATA_H * SLIDE_H_CM)


def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── 色 ───────────────────────────────────────────────────────────────────
C_MAG   = "#3b82f6"
C_MOD   = "#f97316"
C_FILM  = "#8b5cf6"
C_ENC   = "#10b981"
C_HEAD  = "#6b7280"
C_INPUT = "#374151"
C_BG_MAG  = "#eff6ff"
C_BG_MOD  = "#fff7ed"
C_BG_ENC  = "#ecfdf5"
C_BG_HEAD = "#f9fafb"

# ── Presentation 初期化 ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = cm(SLIDE_W_CM)
prs.slide_height = cm(SLIDE_H_CM)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes


# ─────────────────────────────────────────────────────────────────────────
# Helper: 角丸矩形
# (x, y) = data 座標 左下角, (w, h) = data 幅・高さ
# ─────────────────────────────────────────────────────────────────────────
def add_box(x, y, w, h, fill_hex, line_hex, line_pt=1.5, alpha=1.0):
    """alpha: 0=透明, 1=不透明"""
    left   = to_emu_x(x)
    top    = to_emu_y(y + h)     # 上端 = y+h をy軸反転
    width  = to_emu_w(w)
    height = to_emu_h(h)

    shape = shapes.add_shape(5, left, top, width, height)  # 5=rounded rect

    # 塗りつぶし (solidFill + alpha)
    sp = shape.element
    spPr = sp.find(qn("p:spPr"))
    # 既存 fill 要素を削除して solidFill を直接 XML で設定
    for tag in ("noFill", "solidFill", "gradFill", "pattFill"):
        for e in spPr.findall(qn("a:" + tag)):
            spPr.remove(e)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    h_val = fill_hex.lstrip("#")
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", h_val)
    if alpha < 1.0:
        alphaElem = etree.SubElement(srgbClr, qn("a:alpha"))
        alphaElem.set("val", str(int(alpha * 100000)))

    # 枠線
    shape.line.color.rgb = rgb(line_hex)
    shape.line.width = Pt(line_pt)

    return shape


# ─────────────────────────────────────────────────────────────────────────
# Helper: テキストボックス
# center_x, center_y = data 座標中心
# ─────────────────────────────────────────────────────────────────────────
def add_text(cx, cy, w, h, text, font_pt=11, bold=False,
             color_hex="#000000", align=PP_ALIGN.CENTER):
    """
    複数行テキスト: '\\n' で段落分割。
    """
    left   = to_emu_x(cx - w / 2)
    top    = to_emu_y(cy + h / 2)
    width  = to_emu_w(w)
    height = to_emu_h(h)

    txb = shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = False

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_pt)
        run.font.bold = bold
        run.font.color.rgb = rgb(color_hex)

    return txb


# ─────────────────────────────────────────────────────────────────────────
# Helper: 矢印 (直線コネクター + 先端矢印)
# ─────────────────────────────────────────────────────────────────────────
def add_arrow(x1, y1, x2, y2, color_hex="#374151", line_pt=1.5):
    """始点 (x1,y1) → 終点 (x2,y2), data 座標"""
    bx = to_emu_x(x1)
    by = to_emu_y(y1)
    ex = to_emu_x(x2)
    ey = to_emu_y(y2)

    conn = shapes.add_connector(1, bx, by, ex, ey)  # 1=straight

    # 色・太さ
    sp = conn.element
    spPr = sp.find(qn("p:spPr"))

    # <a:ln> を追加（色・太さ・矢印頭）
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(Pt(line_pt))))

    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", color_hex.lstrip("#"))

    # 先端矢印 (headEnd)
    headEnd = etree.SubElement(ln, qn("a:headEnd"))
    headEnd.set("type", "arrow")
    headEnd.set("w", "med")
    headEnd.set("len", "med")

    # pstyle の schemeClr を上書き（PowerPoint のデフォルト色を消す）
    style = sp.find(qn("p:style"))
    if style is not None:
        lnRef = style.find(qn("a:lnRef"))
        if lnRef is not None:
            # schemeClr を削除して srgbClr に差し替え
            for child in list(lnRef):
                lnRef.remove(child)
            srgb2 = etree.SubElement(lnRef, qn("a:srgbClr"))
            srgb2.set("val", color_hex.lstrip("#"))

    return conn


# ═════════════════════════════════════════════════════════════════════════
# 描画開始
# ═════════════════════════════════════════════════════════════════════════

# ─────────────────────────────────────────────────────────────────────────
# Section backgrounds
# ─────────────────────────────────────────────────────────────────────────
add_box(1.8, 1.2, 5.0, 4.6, C_BG_MAG, "#93c5fd", line_pt=1.0, alpha=0.5)
add_text(4.3, 5.60, 3.8, 0.45, "Dual-Stream Embedding",
         font_pt=11, bold=True, color_hex="#1d4ed8")

add_box(8.9, 1.2, 3.2, 4.6, C_BG_ENC, "#6ee7b7", line_pt=1.0, alpha=0.5)
add_text(10.5, 5.60, 3.2, 0.45, "Transformer Encoder",
         font_pt=11, bold=True, color_hex="#065f46")

add_box(12.5, 1.2, 3.1, 4.6, C_BG_HEAD, "#d1d5db", line_pt=1.0, alpha=0.5)
add_text(14.05, 5.57, 3.0, 0.45, "Output Heads",
         font_pt=11, bold=True, color_hex="#374151")

# ─────────────────────────────────────────────────────────────────────────
# Input box
# ─────────────────────────────────────────────────────────────────────────
add_box(0.15, 3.0, 1.5, 1.0, "#f3f4f6", C_INPUT, line_pt=1.8)
add_text(0.90, 3.70, 1.4, 0.40, "Input", font_pt=13, bold=True, color_hex=C_INPUT)
add_text(0.90, 3.28, 1.4, 0.40, "x\u2081, \u2026, x_L", font_pt=12, color_hex=C_INPUT)

# ─────────────────────────────────────────────────────────────────────────
# Feature Extraction
# ─────────────────────────────────────────────────────────────────────────
# Magnitude feature box
add_box(1.95, 4.35, 1.9, 0.80, "#dbeafe", C_MAG, line_pt=1.5)
add_text(2.90, 4.83, 1.8, 0.38, "Mag. Features", font_pt=12, color_hex=C_MAG)
add_text(2.90, 4.50, 1.8, 0.38, "f^mag \u2208 R\u2074",  font_pt=11, color_hex=C_MAG)

# Modulo feature box
add_box(1.95, 2.85, 1.9, 0.80, "#ffedd5", C_MOD, line_pt=1.5)
add_text(2.90, 3.33, 1.8, 0.38, "Mod. Features",           font_pt=12, color_hex=C_MOD)
add_text(2.90, 3.00, 1.8, 0.38, "f^mod \u2208 R\u00b2\u2070\u2070", font_pt=11, color_hex=C_MOD)

# ─────────────────────────────────────────────────────────────────────────
# Projection layers
# ─────────────────────────────────────────────────────────────────────────
# MLP_mag
add_box(4.15, 4.35, 1.5, 0.80, "#bfdbfe", C_MAG, line_pt=1.5)
add_text(4.90, 4.83, 1.4, 0.38, "MLP_mag",          font_pt=13, bold=True, color_hex=C_MAG)
add_text(4.90, 4.50, 1.4, 0.38, "h^mag \u2208 R^d", font_pt=10, color_hex=C_MAG)

# W_mod
add_box(4.15, 2.85, 1.5, 0.80, "#fed7aa", C_MOD, line_pt=1.5)
add_text(4.90, 3.33, 1.4, 0.38, "W_mod",            font_pt=13, bold=True, color_hex=C_MOD)
add_text(4.90, 3.00, 1.4, 0.38, "h^mod \u2208 R^d", font_pt=10, color_hex=C_MOD)

# ─────────────────────────────────────────────────────────────────────────
# FiLM Fusion
# ─────────────────────────────────────────────────────────────────────────
add_box(6.05, 3.55, 2.55, 1.90, "#ede9fe", C_FILM, line_pt=1.8)
add_text(7.40, 5.12, 2.4, 0.38, "FiLM Fusion",
         font_pt=12, bold=True, color_hex=C_FILM)
add_text(7.40, 4.73, 2.4, 0.38, "(1+\u03b3)\u2299h^mag+\u03b2",
         font_pt=11, bold=True, color_hex=C_FILM)
add_text(7.40, 4.37, 2.4, 0.38, "e_i \u2208 R^d",
         font_pt=12, color_hex=C_FILM)
add_text(7.40, 4.02, 2.4, 0.38, "+ PE",
         font_pt=12, color_hex="#7c3aed")

# ─────────────────────────────────────────────────────────────────────────
# Transformer Encoder (stacked boxes for depth effect)
# ─────────────────────────────────────────────────────────────────────────
for i, alpha_v in enumerate([0.70, 0.85, 1.0]):
    yoff = (2 - i) * 0.22
    add_box(9.05 + (2 - i) * 0.12, 3.25 - yoff, 2.6, 1.5,
            "#d1fae5", C_ENC, line_pt=1.2, alpha=alpha_v)

add_text(10.35, 4.27, 2.5, 0.38, "Pre-LN",        font_pt=13, bold=True, color_hex=C_ENC)
add_text(10.35, 3.94, 2.5, 0.38, "Multi-Head Attn", font_pt=12, color_hex=C_ENC)
add_text(10.35, 3.63, 2.5, 0.38, "FFN",            font_pt=13, color_hex=C_ENC)
add_text(10.35, 3.31, 2.5, 0.38, "\u00d7N layers", font_pt=13, bold=True, color_hex="#065f46")

# ─────────────────────────────────────────────────────────────────────────
# Output Heads
# ─────────────────────────────────────────────────────────────────────────
HEAD_X = 12.65
HEAD_W = 2.75
HEAD_SPECS = [
    (4.40, "#fef3c7", "#d97706", "Magnitude\nRegression",       "v^, \u03c3\u00b2^ \u2208 R"),
    (3.05, "#e0f2fe", "#0284c7", "Sign\nClassification",        "s^ \u2208 {+, \u2212, 0}"),
    (1.70, "#fce7f3", "#be185d", "Modulo \u00d7100\nClassification", "r^(m) \u2208 {0,\u2026,m\u22121}"),
]

for (y0, fc, ec, title_text, subtitle_text) in HEAD_SPECS:
    add_box(HEAD_X, y0, HEAD_W, 1.1, fc, ec, line_pt=1.5)
    title_lines = title_text.split("\n")
    # タイトル行1
    add_text(HEAD_X + HEAD_W / 2, y0 + 0.89, HEAD_W - 0.1, 0.35,
             title_lines[0], font_pt=11, bold=True, color_hex=ec)
    # タイトル行2
    if len(title_lines) > 1:
        add_text(HEAD_X + HEAD_W / 2, y0 + 0.60, HEAD_W - 0.1, 0.35,
                 title_lines[1], font_pt=11, bold=True, color_hex=ec)
    # サブタイトル
    add_text(HEAD_X + HEAD_W / 2, y0 + 0.24, HEAD_W - 0.1, 0.35,
             subtitle_text, font_pt=10, color_hex="#374151")

# ─────────────────────────────────────────────────────────────────────────
# Arrows
# ─────────────────────────────────────────────────────────────────────────
# Input → Mag feature (上方向に斜め)
add_arrow(1.65, 3.75, 1.95, 4.75, C_MAG, line_pt=1.4)
# Input → Mod feature (水平)
add_arrow(1.65, 3.25, 1.95, 3.25, C_MOD, line_pt=1.4)

# Mag feature → MLP_mag
add_arrow(3.85, 4.75, 4.15, 4.75, C_MAG, line_pt=1.4)
# Mod feature → W_mod
add_arrow(3.85, 3.25, 4.15, 3.25, C_MOD, line_pt=1.4)

# MLP_mag → FiLM (h_mag)
add_arrow(5.65, 4.75, 6.05, 4.75, C_MAG, line_pt=1.4)
# W_mod → FiLM (h_mod, 斜め上)
add_arrow(5.65, 3.25, 6.05, 4.05, C_MOD, line_pt=1.4)

# FiLM → Encoder
add_arrow(8.60, 4.50, 8.90, 4.00, C_FILM, line_pt=1.6)

# Encoder → Output Heads (fan-out)
for y_head in [4.95, 3.60, 2.25]:
    add_arrow(11.78, 4.00, 12.50, y_head, C_HEAD, line_pt=1.2)

# ─────────────────────────────────────────────────────────────────────────
# Legend (色サンプル + ラベル)
# ─────────────────────────────────────────────────────────────────────────
legend_items = [
    ("#dbeafe", C_MAG,  "Magnitude stream"),
    ("#ffedd5", C_MOD,  "Modulo stream"),
    ("#ede9fe", C_FILM, "FiLM fusion"),
    ("#d1fae5", C_ENC,  "Transformer Encoder"),
    ("#f9fafb", C_HEAD, "Output heads"),
]

leg_x0    = 0.3
leg_y     = 0.52
swatch_w  = 0.45
swatch_h  = 0.38
item_w    = 3.0

for i, (fc, ec, lbl) in enumerate(legend_items):
    lx = leg_x0 + i * item_w
    # 色サンプル
    add_box(lx, leg_y - swatch_h / 2, swatch_w, swatch_h,
            fc, ec, line_pt=1.2)
    # ラベル
    add_text(lx + swatch_w + 0.85, leg_y, item_w - swatch_w - 0.1, 0.40,
             lbl, font_pt=10, color_hex="#374151", align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────
# 保存
# ─────────────────────────────────────────────────────────────────────────
out_path = OUT_DIR / "fig1_architecture.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
